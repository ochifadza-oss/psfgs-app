"""
PSFGS Training Presentation Generator
Generates a professional PowerPoint training deck for the
Public Sector Financial Governance Suite (PSFGS).
Uses python-pptx (v1.0.2).
"""

from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
import copy

# ---------------------------------------------------------------------------
# OUTPUT PATH
# ---------------------------------------------------------------------------
OUTPUT_PATH = r"C:\FY 2025 - 2026\GMT\RSA\APPS_MARKET\MUNICIPALITIES - GOVERNANCE\psfgs-app\PSFGS_Training_Deck.pptx"

# ---------------------------------------------------------------------------
# COLOUR PALETTE
# ---------------------------------------------------------------------------
DARK_BLUE   = RGBColor(31,  56,  100)   # #1F3864
MID_BLUE    = RGBColor(46,  117, 182)   # #2E75B6
LIGHT_BLUE  = RGBColor(238, 245, 252)   # #EEF5FC
DARK_GREEN  = RGBColor(26,  107,  60)   # #1A6B3C
AMBER       = RGBColor(212, 134,  26)   # #D4861A
RED         = RGBColor(192,  57,  43)   # #C0392B
WHITE       = RGBColor(255, 255, 255)
GREY        = RGBColor(90,  100, 120)   # #5A6478
BODY        = RGBColor(44,   62,  80)   # #2C3E50
LIGHT_GREY  = RGBColor(244, 246, 249)   # #F4F6F9
LIGHT_BLUE2 = RGBColor(174, 214, 241)   # #AED6F1  (subtitle accent on title slide)
PALE_GREEN  = RGBColor(213, 232, 212)   # light green fill for boxes
PALE_AMBER  = RGBColor(252, 228, 187)   # light amber fill
PALE_BLUE   = RGBColor(189, 215, 238)   # pale blue fill

# ---------------------------------------------------------------------------
# SLIDE DIMENSIONS  33.87cm × 19.05cm
# ---------------------------------------------------------------------------
SLIDE_W = Cm(33.87)
SLIDE_H = Cm(19.05)

# ---------------------------------------------------------------------------
# HELPER UTILITIES
# ---------------------------------------------------------------------------

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    """Add a completely blank slide (no placeholders)."""
    blank_layout = prs.slide_layouts[6]   # index 6 = blank in most themes
    return prs.slides.add_slide(blank_layout)


def add_rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_width_pt=0):
    """Add a filled rectangle shape."""
    from pptx.util import Pt as _Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    fill = shape.fill
    if fill_rgb:
        fill.solid()
        fill.fore_color.rgb = fill_rgb
    else:
        fill.background()
    line = shape.line
    if line_rgb:
        line.color.rgb = line_rgb
        line.width = Pt(line_width_pt) if line_width_pt else Pt(1)
    else:
        line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text, font_size=12, bold=False, italic=False,
                color=None, align=PP_ALIGN.LEFT, font_name="Calibri",
                wrap=True, v_anchor=None):
    """Add a text box with given properties."""
    from pptx.enum.text import MSO_ANCHOR
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    if v_anchor:
        tf.vertical_anchor = v_anchor
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name      = font_name
    run.font.size      = Pt(font_size)
    run.font.bold      = bold
    run.font.italic    = italic
    if color:
        run.font.color.rgb = color
    return txb


def add_text_in_rect(slide, rect_shape, text, font_size=12, bold=False,
                     italic=False, color=None, align=PP_ALIGN.LEFT,
                     font_name="Calibri"):
    """Overlay a textbox aligned to a rectangle's bounding box."""
    x = rect_shape.left
    y = rect_shape.top
    w = rect_shape.width
    h = rect_shape.height
    return add_textbox(slide, x, y, w, h, text, font_size=font_size,
                       bold=bold, italic=italic, color=color, align=align,
                       font_name=font_name)


def add_para(tf, text, font_size=12, bold=False, italic=False,
             color=None, align=PP_ALIGN.LEFT, font_name="Calibri", space_before=0):
    """Append a paragraph to an existing text frame."""
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.name   = font_name
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return p


def header_band(slide, title_text, sub_text=None, band_h=Cm(2.8)):
    """Dark-blue header band across the top of the slide."""
    add_rect(slide, 0, 0, SLIDE_W, band_h, fill_rgb=DARK_BLUE)
    add_textbox(slide, Cm(0.6), Cm(0.2), SLIDE_W - Cm(1.2), band_h - Cm(0.2),
                title_text, font_size=22, bold=True, color=WHITE,
                align=PP_ALIGN.LEFT, wrap=False)
    if sub_text:
        add_textbox(slide, Cm(0.6), band_h - Cm(0.85), SLIDE_W - Cm(1.2), Cm(0.85),
                    sub_text, font_size=11, bold=False, color=LIGHT_BLUE2,
                    align=PP_ALIGN.LEFT, wrap=False)


def light_bg(slide):
    """Add a light grey background to a slide."""
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=LIGHT_GREY)


def content_box(slide, x, y, w, h, heading, bullets, fill_rgb, border_rgb,
                heading_color=None, bullet_color=None, heading_size=13, bullet_size=11):
    """Draw a coloured content box with heading and bullet list."""
    add_rect(slide, x, y, w, h, fill_rgb=fill_rgb, line_rgb=border_rgb, line_width_pt=1.5)
    # heading
    hh = Cm(0.9)
    add_textbox(slide, x + Cm(0.25), y + Cm(0.15), w - Cm(0.5), hh,
                heading, font_size=heading_size, bold=True,
                color=heading_color or border_rgb, align=PP_ALIGN.LEFT)
    # bullets
    txb = slide.shapes.add_textbox(x + Cm(0.25), y + hh + Cm(0.05),
                                   w - Cm(0.5), h - hh - Cm(0.3))
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = b
        run.font.name  = "Calibri"
        run.font.size  = Pt(bullet_size)
        run.font.color.rgb = bullet_color or BODY


def step_flow(slide, y_top, steps, box_w=None, box_h=Cm(1.45)):
    """Draw a horizontal numbered step flow with arrows."""
    n = len(steps)
    margin = Cm(0.7)
    total_w = SLIDE_W - 2 * margin
    if box_w is None:
        # fit boxes and arrows in available width
        arrow_w = Cm(0.7)
        box_w = (total_w - (n - 1) * arrow_w) / n

    x = margin
    colors = [DARK_BLUE, MID_BLUE, DARK_GREEN, AMBER, GREY, RED, DARK_BLUE]

    for i, step in enumerate(steps):
        col = colors[i % len(colors)]
        add_rect(slide, x, y_top, box_w, box_h, fill_rgb=col)
        # number circle approximation — just bold number prefix in text
        add_textbox(slide, x + Cm(0.15), y_top + Cm(0.15),
                    box_w - Cm(0.3), box_h - Cm(0.3),
                    f"{i+1}. {step}", font_size=9.5, bold=False, color=WHITE,
                    align=PP_ALIGN.LEFT)
        x += box_w
        if i < n - 1:
            # arrow as narrow dark-blue rect
            add_rect(slide, x, y_top + box_h / 2 - Cm(0.15),
                     Cm(0.6), Cm(0.3), fill_rgb=GREY)
            # arrowhead text
            add_textbox(slide, x - Cm(0.05), y_top + box_h / 2 - Cm(0.3),
                        Cm(0.7), Cm(0.6), ">", font_size=11, bold=True,
                        color=GREY, align=PP_ALIGN.CENTER)
            x += Cm(0.7)


def feature_boxes(slide, y_top, features, cols=4):
    """Draw a row of feature boxes."""
    margin = Cm(0.7)
    gap    = Cm(0.35)
    n = len(features)
    total_w = SLIDE_W - 2 * margin
    box_w = (total_w - gap * (cols - 1)) / cols
    box_h = Cm(3.2)

    fill_colors   = [DARK_BLUE,  MID_BLUE,   DARK_GREEN, AMBER,
                     RED,        GREY,        DARK_BLUE,  MID_BLUE]
    border_colors = [MID_BLUE,   DARK_BLUE,  PALE_GREEN, PALE_AMBER,
                     RED,        GREY,        MID_BLUE,   DARK_BLUE]

    for i, (feat_title, feat_desc) in enumerate(features[:cols]):
        col_fill = fill_colors[i % len(fill_colors)]
        x = margin + i * (box_w + gap)
        add_rect(slide, x, y_top, box_w, box_h, fill_rgb=col_fill)
        # Feature title
        add_textbox(slide, x + Cm(0.2), y_top + Cm(0.2),
                    box_w - Cm(0.4), Cm(0.9),
                    feat_title, font_size=11, bold=True, color=WHITE,
                    align=PP_ALIGN.LEFT)
        # Feature description
        add_textbox(slide, x + Cm(0.2), y_top + Cm(1.1),
                    box_w - Cm(0.4), box_h - Cm(1.3),
                    feat_desc, font_size=9.5, bold=False, color=WHITE,
                    align=PP_ALIGN.LEFT)


# ===========================================================================
# SLIDE 1 — TITLE SLIDE
# ===========================================================================
def slide_title(prs):
    slide = blank_slide(prs)
    # Full dark-blue background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=DARK_BLUE)

    # Decorative top bar (thinner MID_BLUE strip)
    add_rect(slide, 0, 0, SLIDE_W, Cm(0.5), fill_rgb=MID_BLUE)

    # "PSFGS" main title
    add_textbox(slide, Cm(1), Cm(2.8), SLIDE_W - Cm(2), Cm(3.2),
                "PSFGS", font_size=54, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, Cm(1), Cm(5.8), SLIDE_W - Cm(2), Cm(1.2),
                "Public Sector Financial Governance Suite",
                font_size=16, bold=False, color=LIGHT_BLUE2,
                align=PP_ALIGN.CENTER)

    # Horizontal white divider line
    add_rect(slide, Cm(4), Cm(7.4), SLIDE_W - Cm(8), Cm(0.08), fill_rgb=WHITE)

    # "USER TRAINING MANUAL"
    add_textbox(slide, Cm(1), Cm(7.7), SLIDE_W - Cm(2), Cm(1.6),
                "USER TRAINING MANUAL",
                font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Tagline
    add_textbox(slide, Cm(1), Cm(9.5), SLIDE_W - Cm(2), Cm(1.0),
                "Complete Guide to All 13 Modules",
                font_size=14, bold=False, italic=True, color=LIGHT_BLUE2,
                align=PP_ALIGN.CENTER)

    # Decorative bottom strip
    add_rect(slide, 0, SLIDE_H - Cm(1.5), SLIDE_W, Cm(1.5), fill_rgb=MID_BLUE)

    # Footer text
    add_textbox(slide, Cm(1), SLIDE_H - Cm(1.3), SLIDE_W - Cm(2), Cm(1.0),
                "Version 1.0  |  April 2026  |  Glance Management Technologies (Pty) Ltd",
                font_size=9, bold=False, color=WHITE, align=PP_ALIGN.CENTER)


# ===========================================================================
# SLIDE 2 — TRAINING AGENDA
# ===========================================================================
def slide_agenda(prs):
    slide = blank_slide(prs)
    light_bg(slide)
    header_band(slide, "Training Agenda", "What We Will Cover Today")

    # Two column backgrounds
    col_w = (SLIDE_W - Cm(2.1)) / 2
    add_rect(slide, Cm(0.7), Cm(3.1), col_w, Cm(13.5), fill_rgb=WHITE,
             line_rgb=PALE_BLUE, line_width_pt=1)
    add_rect(slide, Cm(0.7) + col_w + Cm(0.35), Cm(3.1), col_w, Cm(13.5),
             fill_rgb=WHITE, line_rgb=PALE_BLUE, line_width_pt=1)

    # Left column heading
    add_textbox(slide, Cm(1.0), Cm(3.2), col_w - Cm(0.3), Cm(0.7),
                "Modules 1 - 7", font_size=12, bold=True, color=DARK_BLUE)

    left_modules = [
        "01  Executive Dashboard (DASH)",
        "02  Audit Findings Tracker (AFT)",
        "03  Asset Verification System (AVS)",
        "04  Budget vs Actual Monitor (BVM)",
        "05  IFWE Register (IFW)",
        "06  SCM Compliance Checker (SCC)",
        "07  Consequence Management Tracker (CMT)",
    ]
    txb = slide.shapes.add_textbox(Cm(1.0), Cm(4.0), col_w - Cm(0.3), Cm(11.5))
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for m in left_modules:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = m
        run.font.name  = "Calibri"
        run.font.size  = Pt(12)
        run.font.color.rgb = BODY

    # Right column heading
    rx = Cm(0.7) + col_w + Cm(0.35)
    add_textbox(slide, rx + Cm(0.3), Cm(3.2), col_w - Cm(0.3), Cm(0.7),
                "Modules 8 - 13 & Appendices", font_size=12, bold=True, color=DARK_BLUE)

    right_modules = [
        "08  Delegation of Authority Register (DAR)",
        "09  Policy Compliance Manager (PCR)",
        "10  Section 71 Auto-Generator (S71)",
        "11  Performance KPI Dashboard (PKD)",
        "12  Executive Reporting Portal (ERP)",
        "13  Entity Settings (SET)",
        "",
        "Appendix A  -  User Roles & Permissions Matrix",
        "Appendix B  -  Glossary of Terms",
        "Appendix C  -  Support & Next Steps",
    ]
    txb2 = slide.shapes.add_textbox(rx + Cm(0.3), Cm(4.0), col_w - Cm(0.3), Cm(11.5))
    tf2  = txb2.text_frame
    tf2.word_wrap = True
    first = True
    for m in right_modules:
        if first:
            p = tf2.paragraphs[0]; first = False
        else:
            p = tf2.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = m
        run.font.name  = "Calibri"
        run.font.size  = Pt(11.5 if m.startswith("Appendix") else 12)
        run.font.italic = m.startswith("Appendix")
        run.font.color.rgb = GREY if m.startswith("Appendix") else BODY

    # Bottom note
    add_rect(slide, Cm(0.7), SLIDE_H - Cm(1.35), SLIDE_W - Cm(1.4), Cm(0.95),
             fill_rgb=LIGHT_BLUE, line_rgb=MID_BLUE, line_width_pt=1)
    add_textbox(slide, Cm(1.0), SLIDE_H - Cm(1.3), SLIDE_W - Cm(2.0), Cm(0.85),
                "Estimated training duration: Full day (6 hours) or self-paced online",
                font_size=10, bold=False, italic=True, color=DARK_BLUE,
                align=PP_ALIGN.CENTER)


# ===========================================================================
# SLIDE 3 — SYSTEM OVERVIEW
# ===========================================================================
def slide_overview(prs):
    slide = blank_slide(prs)
    light_bg(slide)
    header_band(slide, "About PSFGS", "System Overview")

    box_w  = (SLIDE_W - Cm(2.8)) / 3
    box_h  = Cm(12.5)
    y_top  = Cm(3.2)
    gap    = Cm(0.35)

    boxes = [
        {
            "heading": "What is PSFGS?",
            "fill":    LIGHT_BLUE,
            "border":  MID_BLUE,
            "hcolor":  DARK_BLUE,
            "bullets": [
                "- Cloud-based SaaS platform built specifically for South African municipalities",
                "- 13 fully integrated governance and financial management modules",
                "- MFMA-native design — every feature maps to a legislative requirement",
                "- Multi-tenant architecture — one system serves many municipalities",
                "- Role-based access control ensures the right people see the right data",
            ],
        },
        {
            "heading": "Why was it built?",
            "fill":    RGBColor(213, 232, 212),
            "border":  DARK_GREEN,
            "hcolor":  DARK_GREEN,
            "bullets": [
                "- Manual processes are the root cause of most audit failures",
                "- 163+ municipalities received adverse or disclaimed opinions in 2023/24",
                "- The same AGSA findings repeat year after year",
                "- Municipalities lack integrated tools to track, act and report in real time",
                "- PSFGS closes the gap between policy intent and implementation reality",
            ],
        },
        {
            "heading": "Who is it for?",
            "fill":    RGBColor(252, 228, 187),
            "border":  AMBER,
            "hcolor":  RGBColor(160, 100, 10),
            "bullets": [
                "- All municipalities governed by the MFMA",
                "- Metropolitan, District and Local municipalities",
                "- Scales from 8 to 2,000+ concurrent users",
                "- Deployable as a single module or the full 13-module suite",
                "- Designed for CFOs, Accounting Officers, Directors and operational staff",
            ],
        },
    ]

    for i, box in enumerate(boxes):
        x = Cm(0.7) + i * (box_w + gap)
        content_box(slide, x, y_top, box_w, box_h,
                    box["heading"], box["bullets"],
                    fill_rgb=box["fill"], border_rgb=box["border"],
                    heading_color=box["hcolor"], bullet_color=BODY,
                    heading_size=13, bullet_size=10.5)


# ===========================================================================
# SLIDE 4 — LOGGING IN
# ===========================================================================
def slide_login(prs):
    slide = blank_slide(prs)
    light_bg(slide)
    header_band(slide, "Getting Started", "System Login")

    steps = [
        ("Step 1 — Open Browser",
         "Navigate to your municipality's PSFGS URL\n(provided by your System Administrator)"),
        ("Step 2 — Enter Email",
         "Type your email address\nExample: cfo@buffalocity.gov.za"),
        ("Step 3 — Enter Password",
         "Default password: Admin@2026!\nYou will be prompted to change this on first login"),
        ("Step 4 — Sign In",
         "Click the 'Sign In' button\nYou are directed to your role-specific dashboard"),
    ]

    box_h = Cm(3.6)
    box_w = (SLIDE_W - Cm(2.1) - Cm(1.05)) / 4
    y_top = Cm(3.4)

    fills = [DARK_BLUE, MID_BLUE, DARK_GREEN, AMBER]
    for i, (title, body) in enumerate(steps):
        x = Cm(0.7) + i * (box_w + Cm(0.35))
        add_rect(slide, x, y_top, box_w, box_h, fill_rgb=fills[i])
        # Number badge
        add_rect(slide, x + Cm(0.2), y_top + Cm(0.2), Cm(0.85), Cm(0.85),
                 fill_rgb=WHITE)
        add_textbox(slide, x + Cm(0.2), y_top + Cm(0.15), Cm(0.85), Cm(0.85),
                    str(i + 1), font_size=13, bold=True, color=fills[i],
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Cm(0.2), y_top + Cm(1.15), box_w - Cm(0.4), Cm(1.0),
                    title, font_size=11, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
        add_textbox(slide, x + Cm(0.2), y_top + Cm(2.1), box_w - Cm(0.4), Cm(1.3),
                    body, font_size=9.5, bold=False, color=WHITE, align=PP_ALIGN.LEFT)

    # Info callout
    y_info = y_top + box_h + Cm(0.6)
    add_rect(slide, Cm(0.7), y_info, SLIDE_W - Cm(1.4), Cm(2.2),
             fill_rgb=LIGHT_BLUE, line_rgb=MID_BLUE, line_width_pt=1.5)
    add_textbox(slide, Cm(1.0), y_info + Cm(0.2), SLIDE_W - Cm(2.0), Cm(0.7),
                "Information", font_size=12, bold=True, color=DARK_BLUE)
    add_textbox(slide, Cm(1.0), y_info + Cm(0.8), SLIDE_W - Cm(2.0), Cm(1.2),
                "The login page displays your municipality's logo and name. "
                "If these details are incorrect, contact your System Administrator "
                "to update the Entity Settings (Module 13).",
                font_size=10.5, bold=False, color=BODY)

    # Bottom tip
    add_rect(slide, Cm(0.7), SLIDE_H - Cm(1.55), SLIDE_W - Cm(1.4), Cm(1.15),
             fill_rgb=DARK_BLUE)
    add_textbox(slide, Cm(1.0), SLIDE_H - Cm(1.45), SLIDE_W - Cm(2.0), Cm(1.0),
                "Security Tip: Never share your password. Use a strong password of at least 12 characters with letters, numbers and special characters.",
                font_size=9.5, bold=False, italic=True, color=WHITE,
                align=PP_ALIGN.LEFT)


# ===========================================================================
# SLIDE 5 — NAVIGATION OVERVIEW
# ===========================================================================
def slide_navigation(prs):
    slide = blank_slide(prs)
    light_bg(slide)
    header_band(slide, "Navigation Overview", "Finding Your Way Around")

    # Sidebar mockup
    sb_w = Cm(7.5)
    sb_h = SLIDE_H - Cm(3.0)
    sb_x = Cm(0.7)
    sb_y = Cm(3.1)
    add_rect(slide, sb_x, sb_y, sb_w, sb_h, fill_rgb=DARK_BLUE)
    add_textbox(slide, sb_x + Cm(0.3), sb_y + Cm(0.25), sb_w - Cm(0.6), Cm(0.7),
                "PSFGS Menu", font_size=12, bold=True, color=WHITE)

    sidebar_items = [
        ("OVERVIEW",         [" Dashboard", " Finding Workflow"]),
        ("DOMAIN A - Audit", [" AFT  Audit Findings Tracker", " AVS  Asset Verification"]),
        ("DOMAIN B - Finance",[" BVM  Budget vs Actual", " IFW  IFWE Register"]),
        ("DOMAIN C - Report", [" ERP  Exec Reporting Portal", " S71  Section 71 Report", " PKD  Performance KPIs"]),
        ("DOMAIN D - Comply", [" SCC  SCM Compliance", " CMT  Consequence Mgmt", " DAR  Delegations", " PCR  Policy Compliance"]),
        ("SETTINGS",         [" Entity Settings", " User Management"]),
    ]

    y_cur = sb_y + Cm(1.1)
    for section, items in sidebar_items:
        add_rect(slide, sb_x, y_cur, sb_w, Cm(0.55), fill_rgb=MID_BLUE)
        add_textbox(slide, sb_x + Cm(0.3), y_cur + Cm(0.05), sb_w - Cm(0.4), Cm(0.5),
                    section, font_size=8.5, bold=True, color=WHITE)
        y_cur += Cm(0.55)
        for item in items:
            add_textbox(slide, sb_x + Cm(0.3), y_cur, sb_w - Cm(0.4), Cm(0.5),
                        item, font_size=8.5, bold=False, color=LIGHT_BLUE2)
            y_cur += Cm(0.52)
        y_cur += Cm(0.15)

    # Right side descriptions
    rx = sb_x + sb_w + Cm(0.7)
    rw = SLIDE_W - rx - Cm(0.7)

    add_textbox(slide, rx, Cm(3.2), rw, Cm(0.7),
                "How to Navigate the System", font_size=14, bold=True, color=DARK_BLUE)

    right_sections = [
        ("Overview", MID_BLUE, "The default landing area after login. Shows the Executive Dashboard with live KPI tiles and a diagram of the finding workflow across all modules."),
        ("Domain A — Audit Readiness", DARK_BLUE, "AFT: Manage all audit findings end-to-end. AVS: Conduct and record asset verification sessions in the field."),
        ("Domain B — Financial Control", DARK_GREEN, "BVM: Monitor budget votes against actual spend in real time. IFW: Capture and track all irregular, fruitless and wasteful expenditure."),
        ("Domain C — Reporting", AMBER, "ERP: Executive governance dashboard for Mayor and Council. S71: Auto-generate monthly National Treasury reports. PKD: Track SDBIP KPIs by quarter."),
        ("Domain D — Compliance", MID_BLUE, "SCC: Run automated SCM compliance checks on requisitions. CMT: Track misconduct and disciplinary cases. DAR: Manage Council delegations. PCR: Monitor policy review schedules."),
    ]

    y_sec = Cm(3.95)
    for title, color, desc in right_sections:
        add_rect(slide, rx, y_sec, Cm(0.35), Cm(1.5), fill_rgb=color)
        add_textbox(slide, rx + Cm(0.5), y_sec, rw - Cm(0.5), Cm(0.6),
                    title, font_size=11, bold=True, color=color)
        add_textbox(slide, rx + Cm(0.5), y_sec + Cm(0.6), rw - Cm(0.5), Cm(0.85),
                    desc, font_size=9, bold=False, color=BODY)
        y_sec += Cm(1.65)


# ===========================================================================
# SLIDE 6 — USER ROLES
# ===========================================================================
def slide_user_roles(prs):
    slide = blank_slide(prs)
    light_bg(slide)
    header_band(slide, "User Roles & Access Levels", "Who can do what in PSFGS")

    rows_data = [
        ["Role",             "Who Uses It",                  "Key Modules",                "Access Level"],
        ["System_Admin",     "IT / System Administrator",    "All modules",                "Full — configure everything"],
        ["Accounting_Officer","City Manager / Municipal Mgr","AFT, CMT, DAR, Dashboard",   "Full write across modules"],
        ["CFO",              "Chief Financial Officer",      "BVM, S71, IFW, Dashboard",   "Full write — financial modules"],
        ["Director",         "Directorate Head",             "AFT, AVS, PKD",              "Write within own directorate"],
        ["Manager",          "Senior Manager",               "AFT, AVS, SCC",              "Write within own area"],
        ["Officer",          "Operational Staff",            "AFT (evidence), AVS",        "Limited write — own tasks"],
        ["Auditor",          "Internal Audit staff",         "AFT (all stages), PCR",      "Audit write — cross-directorate"],
        ["Councillor",       "Elected Councillor",           "Dashboard, ERP",             "Read only — view reports"],
    ]

    tbl_x = Cm(0.7)
    tbl_y = Cm(3.2)
    tbl_w = SLIDE_W - Cm(1.4)
    tbl_h = SLIDE_H - Cm(4.2)

    col_widths = [Cm(5.0), Cm(7.5), Cm(9.5), Cm(10.0)]
    row_h = tbl_h / len(rows_data)

    for r_idx, row in enumerate(rows_data):
        y = tbl_y + r_idx * row_h
        x_cur = tbl_x
        for c_idx, cell_text in enumerate(row):
            cw = col_widths[c_idx]
            if r_idx == 0:
                fill = DARK_BLUE
                fc   = WHITE
                fb   = True
            elif r_idx % 2 == 0:
                fill = LIGHT_GREY
                fc   = BODY
                fb   = False
            else:
                fill = WHITE
                fc   = BODY
                fb   = False
            add_rect(slide, x_cur, y, cw, row_h,
                     fill_rgb=fill, line_rgb=PALE_BLUE, line_width_pt=0.5)
            add_textbox(slide, x_cur + Cm(0.2), y + Cm(0.1),
                        cw - Cm(0.4), row_h - Cm(0.1),
                        cell_text, font_size=10, bold=fb, color=fc,
                        align=PP_ALIGN.LEFT)
            x_cur += cw


# ===========================================================================
# MODULE SLIDE FACTORY
# ===========================================================================

def module_intro(prs, mod_num, mod_name, mod_code, purpose_text,
                 legislation, users_list, key_benefits):
    """Slide A — Module Introduction."""
    slide = blank_slide(prs)
    light_bg(slide)

    # Full header band (taller for module slides)
    band_h = Cm(3.3)
    add_rect(slide, 0, 0, SLIDE_W, band_h, fill_rgb=DARK_BLUE)

    # Module badge
    badge_w = Cm(2.2)
    badge_h = Cm(0.85)
    add_rect(slide, Cm(0.7), Cm(0.55), badge_w, badge_h, fill_rgb=MID_BLUE)
    add_textbox(slide, Cm(0.7), Cm(0.55), badge_w, badge_h,
                f"MODULE {mod_num:02d}", font_size=9, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)

    # Code badge
    cb_x = Cm(0.7) + badge_w + Cm(0.3)
    add_rect(slide, cb_x, Cm(0.55), Cm(1.5), badge_h, fill_rgb=AMBER)
    add_textbox(slide, cb_x, Cm(0.55), Cm(1.5), badge_h,
                mod_code, font_size=9, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)

    # Module name
    add_textbox(slide, Cm(0.7), Cm(1.55), SLIDE_W - Cm(1.4), Cm(1.5),
                mod_name, font_size=22, bold=True, color=WHITE,
                align=PP_ALIGN.LEFT)

    # Left content area (60%)
    lw = SLIDE_W * 0.58
    lx = Cm(0.7)
    ly = band_h + Cm(0.35)
    lh = SLIDE_H - ly - Cm(0.4)

    # Purpose box
    ph = Cm(3.5)
    add_rect(slide, lx, ly, lw, ph, fill_rgb=WHITE, line_rgb=PALE_BLUE, line_width_pt=1)
    add_textbox(slide, lx + Cm(0.25), ly + Cm(0.15), lw - Cm(0.5), Cm(0.65),
                "Purpose", font_size=12, bold=True, color=DARK_BLUE)
    add_textbox(slide, lx + Cm(0.25), ly + Cm(0.8), lw - Cm(0.5), ph - Cm(1.0),
                purpose_text, font_size=10.5, bold=False, color=BODY)

    # Legislation box
    legy = ly + ph + Cm(0.3)
    legh = Cm(2.1)
    add_rect(slide, lx, legy, lw, legh, fill_rgb=LIGHT_BLUE, line_rgb=MID_BLUE, line_width_pt=1)
    add_textbox(slide, lx + Cm(0.25), legy + Cm(0.15), lw - Cm(0.5), Cm(0.6),
                "Legislative Basis", font_size=12, bold=True, color=DARK_BLUE)
    add_textbox(slide, lx + Cm(0.25), legy + Cm(0.7), lw - Cm(0.5), legh - Cm(0.8),
                legislation, font_size=10, bold=False, color=BODY)

    # Users box
    uy = legy + legh + Cm(0.3)
    uh = lh - (uy - ly)
    add_rect(slide, lx, uy, lw, uh, fill_rgb=LIGHT_GREY, line_rgb=GREY, line_width_pt=1)
    add_textbox(slide, lx + Cm(0.25), uy + Cm(0.15), lw - Cm(0.5), Cm(0.6),
                "Who Uses This Module", font_size=12, bold=True, color=DARK_BLUE)
    users_text = "   ".join(f"[ {u} ]" for u in users_list)
    add_textbox(slide, lx + Cm(0.25), uy + Cm(0.75), lw - Cm(0.5), uh - Cm(0.9),
                users_text, font_size=10, bold=False, color=MID_BLUE)

    # Right content area — Key Benefits (40%)
    rx = lx + lw + Cm(0.35)
    rw = SLIDE_W - rx - Cm(0.7)
    add_rect(slide, rx, ly, rw, lh, fill_rgb=RGBColor(213, 232, 212),
             line_rgb=DARK_GREEN, line_width_pt=2)
    add_rect(slide, rx, ly, rw, Cm(0.85), fill_rgb=DARK_GREEN)
    add_textbox(slide, rx + Cm(0.25), ly + Cm(0.1), rw - Cm(0.5), Cm(0.7),
                "Key Benefits", font_size=12, bold=True, color=WHITE)

    txb = slide.shapes.add_textbox(rx + Cm(0.25), ly + Cm(1.0),
                                   rw - Cm(0.5), lh - Cm(1.2))
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for b in key_benefits:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = f"+ {b}"
        run.font.name  = "Calibri"
        run.font.size  = Pt(10.5)
        run.font.color.rgb = BODY


def module_features(prs, mod_name, features_4, workflow_steps):
    """Slide B — Key Features & Workflow."""
    slide = blank_slide(prs)
    light_bg(slide)
    header_band(slide, f"{mod_name} — How It Works", "Features & Workflow")

    # 4 feature boxes
    feature_boxes(slide, Cm(3.4), features_4[:4], cols=4)

    # Workflow label
    add_textbox(slide, Cm(0.7), Cm(7.3), SLIDE_W - Cm(1.4), Cm(0.7),
                "Step-by-Step Workflow", font_size=13, bold=True, color=DARK_BLUE)

    # Step flow
    step_flow(slide, Cm(8.0), workflow_steps)

    # Bottom info strip
    y_info = Cm(10.2)
    add_rect(slide, Cm(0.7), y_info, SLIDE_W - Cm(1.4), SLIDE_H - y_info - Cm(0.5),
             fill_rgb=WHITE, line_rgb=PALE_BLUE, line_width_pt=1)
    add_textbox(slide, Cm(1.0), y_info + Cm(0.2),
                SLIDE_W - Cm(2.0), SLIDE_H - y_info - Cm(0.7),
                "Access this module from the left navigation menu after logging in. "
                "All actions are logged with the user's name, date and time for audit trail purposes.",
                font_size=10, bold=False, italic=True, color=GREY)


def module_exercise(prs, mod_name, scenario_text, task_steps,
                    expected_outcome, test_email, test_password="Admin@2026!"):
    """Slide C — Practical Exercise."""
    slide = blank_slide(prs)
    light_bg(slide)
    header_band(slide, f"Hands-On Practice — {mod_name}", "Practical Exercise")

    content_y = Cm(3.25)
    content_h  = SLIDE_H - content_y - Cm(0.5)
    col_w = (SLIDE_W - Cm(2.1) - Cm(0.7)) / 2
    lx = Cm(0.7)
    rx = lx + col_w + Cm(0.35)

    # ---- LEFT COLUMN ----
    # Scenario box (amber)
    sc_h = Cm(4.2)
    add_rect(slide, lx, content_y, col_w, sc_h,
             fill_rgb=RGBColor(252, 228, 187), line_rgb=AMBER, line_width_pt=2)
    add_rect(slide, lx, content_y, col_w, Cm(0.8), fill_rgb=AMBER)
    add_textbox(slide, lx + Cm(0.25), content_y + Cm(0.1), col_w - Cm(0.5), Cm(0.65),
                "Scenario", font_size=12, bold=True, color=WHITE)
    add_textbox(slide, lx + Cm(0.25), content_y + Cm(0.9), col_w - Cm(0.5), sc_h - Cm(1.0),
                scenario_text, font_size=10, bold=False, color=BODY)

    # Your Task box
    task_y = content_y + sc_h + Cm(0.35)
    task_h = content_h - sc_h - Cm(0.35)
    add_rect(slide, lx, task_y, col_w, task_h,
             fill_rgb=WHITE, line_rgb=MID_BLUE, line_width_pt=1.5)
    add_rect(slide, lx, task_y, col_w, Cm(0.8), fill_rgb=MID_BLUE)
    add_textbox(slide, lx + Cm(0.25), task_y + Cm(0.1), col_w - Cm(0.5), Cm(0.65),
                "Your Task", font_size=12, bold=True, color=WHITE)

    txb = slide.shapes.add_textbox(lx + Cm(0.25), task_y + Cm(0.9),
                                   col_w - Cm(0.5), task_h - Cm(1.1))
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for i, step in enumerate(task_steps, 1):
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = f"{i}. {step}"
        run.font.name  = "Calibri"
        run.font.size  = Pt(10.5)
        run.font.color.rgb = BODY

    # ---- RIGHT COLUMN ----
    # Expected outcome (green)
    eo_h = Cm(5.5)
    add_rect(slide, rx, content_y, col_w, eo_h,
             fill_rgb=RGBColor(213, 232, 212), line_rgb=DARK_GREEN, line_width_pt=2)
    add_rect(slide, rx, content_y, col_w, Cm(0.8), fill_rgb=DARK_GREEN)
    add_textbox(slide, rx + Cm(0.25), content_y + Cm(0.1), col_w - Cm(0.5), Cm(0.65),
                "Expected Outcome", font_size=12, bold=True, color=WHITE)
    add_textbox(slide, rx + Cm(0.25), content_y + Cm(0.9), col_w - Cm(0.5), eo_h - Cm(1.0),
                expected_outcome, font_size=10.5, bold=False, color=BODY)

    # Test Login box (blue)
    tl_y = content_y + eo_h + Cm(0.35)
    tl_h = content_h - eo_h - Cm(0.35)
    add_rect(slide, rx, tl_y, col_w, tl_h,
             fill_rgb=LIGHT_BLUE, line_rgb=MID_BLUE, line_width_pt=1.5)
    add_rect(slide, rx, tl_y, col_w, Cm(0.8), fill_rgb=DARK_BLUE)
    add_textbox(slide, rx + Cm(0.25), tl_y + Cm(0.1), col_w - Cm(0.5), Cm(0.65),
                "Test Login Credentials", font_size=12, bold=True, color=WHITE)
    add_textbox(slide, rx + Cm(0.25), tl_y + Cm(0.95), col_w - Cm(0.5), Cm(0.65),
                "Email Address:", font_size=10.5, bold=True, color=DARK_BLUE)
    add_textbox(slide, rx + Cm(0.25), tl_y + Cm(1.55), col_w - Cm(0.5), Cm(0.65),
                test_email, font_size=10.5, bold=False, color=BODY)
    add_textbox(slide, rx + Cm(0.25), tl_y + Cm(2.2), col_w - Cm(0.5), Cm(0.65),
                "Password:", font_size=10.5, bold=True, color=DARK_BLUE)
    add_textbox(slide, rx + Cm(0.25), tl_y + Cm(2.8), col_w - Cm(0.5), Cm(0.65),
                test_password, font_size=10.5, bold=False, color=BODY)
    if tl_h > Cm(4.0):
        add_textbox(slide, rx + Cm(0.25), tl_y + Cm(3.5), col_w - Cm(0.5), tl_h - Cm(3.7),
                    "Note: Change your password after first login.",
                    font_size=9, bold=False, italic=True, color=GREY)


# ===========================================================================
# APPENDIX SLIDES
# ===========================================================================

def slide_roles_matrix(prs):
    slide = blank_slide(prs)
    light_bg(slide)
    header_band(slide, "Appendix A — User Roles & Permissions Matrix",
                "Full permissions across all 13 modules")

    modules_short = ["DASH","AFT","AVS","BVM","IFW","SCC","CMT","DAR","PCR","S71","PKD","ERP","SET"]
    roles = ["Sys_Admin","Acctg_Officer","CFO","Director","Manager","Officer","Auditor","Councillor"]

    # R = Read, W = Write, F = Full, - = None
    perms = {
        "Sys_Admin":       ["F","F","F","F","F","F","F","F","F","F","F","F","F"],
        "Acctg_Officer":   ["F","F","R","R","R","R","F","F","R","R","R","F","-"],
        "CFO":             ["F","R","R","F","F","R","R","R","R","F","R","F","-"],
        "Director":        ["R","W","W","R","R","R","R","R","R","R","W","R","-"],
        "Manager":         ["R","W","W","R","R","W","-","-","R","R","R","R","-"],
        "Officer":         ["-","W","-","-","-","-","-","-","-","-","-","-","-"],
        "Auditor":         ["R","F","R","R","R","R","R","R","W","R","R","R","-"],
        "Councillor":      ["R","-","-","-","-","-","-","-","-","R","R","R","-"],
    }

    tbl_x = Cm(0.7)
    tbl_y = Cm(3.3)
    tbl_w = SLIDE_W - Cm(1.4)
    tbl_h = SLIDE_H - Cm(4.1)

    role_col_w = Cm(4.2)
    mod_col_w  = (tbl_w - role_col_w) / len(modules_short)
    row_h      = tbl_h / (len(roles) + 1)

    # Header row
    add_rect(slide, tbl_x, tbl_y, role_col_w, row_h, fill_rgb=DARK_BLUE,
             line_rgb=MID_BLUE, line_width_pt=0.5)
    add_textbox(slide, tbl_x + Cm(0.15), tbl_y + Cm(0.1), role_col_w - Cm(0.3), row_h - Cm(0.1),
                "Role", font_size=9, bold=True, color=WHITE)
    for c_i, mod in enumerate(modules_short):
        cx = tbl_x + role_col_w + c_i * mod_col_w
        add_rect(slide, cx, tbl_y, mod_col_w, row_h, fill_rgb=DARK_BLUE,
                 line_rgb=MID_BLUE, line_width_pt=0.5)
        add_textbox(slide, cx + Cm(0.05), tbl_y + Cm(0.1), mod_col_w - Cm(0.1), row_h - Cm(0.1),
                    mod, font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    perm_colors = {"F": DARK_GREEN, "W": MID_BLUE, "R": GREY, "-": RGBColor(200,200,200)}

    for r_i, role in enumerate(roles):
        ry = tbl_y + (r_i + 1) * row_h
        bg = LIGHT_GREY if r_i % 2 == 0 else WHITE
        add_rect(slide, tbl_x, ry, role_col_w, row_h, fill_rgb=bg,
                 line_rgb=PALE_BLUE, line_width_pt=0.5)
        add_textbox(slide, tbl_x + Cm(0.15), ry + Cm(0.1), role_col_w - Cm(0.3), row_h - Cm(0.1),
                    role, font_size=9, bold=True, color=DARK_BLUE)
        for c_i, perm in enumerate(perms[role]):
            cx = tbl_x + role_col_w + c_i * mod_col_w
            p_col = perm_colors.get(perm, GREY)
            add_rect(slide, cx, ry, mod_col_w, row_h, fill_rgb=bg,
                     line_rgb=PALE_BLUE, line_width_pt=0.5)
            add_textbox(slide, cx + Cm(0.05), ry + Cm(0.1), mod_col_w - Cm(0.1), row_h - Cm(0.1),
                        perm, font_size=9, bold=(perm == "F"), color=p_col,
                        align=PP_ALIGN.CENTER)

    # Legend
    ly = SLIDE_H - Cm(1.2)
    legend = [("F = Full access", DARK_GREEN), ("W = Write access", MID_BLUE),
              ("R = Read only", GREY), ("- = No access", RGBColor(180,180,180))]
    lx_cur = Cm(0.7)
    for txt, col in legend:
        add_rect(slide, lx_cur, ly, Cm(0.4), Cm(0.4), fill_rgb=col)
        add_textbox(slide, lx_cur + Cm(0.5), ly, Cm(4.5), Cm(0.5),
                    txt, font_size=9, bold=False, color=BODY)
        lx_cur += Cm(5.5)


def slide_glossary(prs):
    slide = blank_slide(prs)
    light_bg(slide)
    header_band(slide, "Appendix B — Glossary of Terms", "Key acronyms and definitions used in PSFGS")

    terms = [
        ("AGSA",     "Auditor-General of South Africa — the constitutional body that audits government finances"),
        ("MFMA",     "Municipal Finance Management Act (Act 56 of 2003) — the primary law governing municipal finances"),
        ("IFWE",     "Irregular, Fruitless and Wasteful Expenditure — expenditure that does not comply with legislation"),
        ("SDBIP",    "Service Delivery and Budget Implementation Plan — quarterly performance plan linked to the IDP and budget"),
        ("IDP",      "Integrated Development Plan — the municipality's five-year development strategy"),
        ("SCM",      "Supply Chain Management — the procurement and asset disposal process"),
        ("PPPFA",    "Preferential Procurement Policy Framework Act — legislation governing B-BBEE scoring in procurement"),
        ("GRAP 17",  "Generally Recognised Accounting Practice 17 — the accounting standard for property, plant and equipment"),
        ("MPAC",     "Municipal Public Accounts Committee — the council committee that holds the executive accountable"),
        ("AFS",      "Annual Financial Statements — the audited financial statements submitted annually"),
        ("NT",       "National Treasury — the national department that oversees municipal financial management"),
        ("COGTA",    "Department of Cooperative Governance and Traditional Affairs — the oversight department for municipalities"),
        ("BVM",      "Budget vs Actual Monitor — PSFGS module for real-time budget variance tracking"),
        ("KPI",      "Key Performance Indicator — a measurable value used to evaluate the achievement of objectives"),
        ("DAR",      "Delegation of Authority Register — PSFGS module for managing Council and MM delegations"),
        ("PCR",      "Policy Compliance Register — PSFGS module for tracking municipal policy review schedules"),
    ]

    col_w = (SLIDE_W - Cm(2.1) - Cm(0.35)) / 2
    lx = Cm(0.7)
    rx = lx + col_w + Cm(0.35)
    y_start = Cm(3.3)
    item_h  = (SLIDE_H - y_start - Cm(0.5)) / 8

    for i, (term, defn) in enumerate(terms):
        col_x = lx if i < 8 else rx
        row_i = i % 8
        y = y_start + row_i * item_h
        bg = LIGHT_GREY if row_i % 2 == 0 else WHITE
        add_rect(slide, col_x, y, col_w, item_h, fill_rgb=bg,
                 line_rgb=PALE_BLUE, line_width_pt=0.5)
        add_rect(slide, col_x, y, Cm(2.3), item_h, fill_rgb=MID_BLUE if row_i % 2 == 0 else DARK_BLUE)
        add_textbox(slide, col_x + Cm(0.1), y + Cm(0.07), Cm(2.1), item_h - Cm(0.1),
                    term, font_size=9.5, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
        add_textbox(slide, col_x + Cm(2.4), y + Cm(0.07), col_w - Cm(2.5), item_h - Cm(0.1),
                    defn, font_size=9, bold=False, color=BODY)


def slide_support(prs):
    slide = blank_slide(prs)
    light_bg(slide)
    header_band(slide, "Appendix C — Support & Next Steps", "Getting help and continuing your learning")

    boxes = [
        {
            "title":  "Getting Help",
            "fill":   LIGHT_BLUE,
            "border": MID_BLUE,
            "tcolor": DARK_BLUE,
            "items":  [
                "Contact your municipality's designated System Administrator for:",
                "- Password resets and account lockouts",
                "- New user account creation",
                "- Module access and permission changes",
                "- Data entry corrections",
                "Your System Admin is listed in the Entity Settings module.",
            ],
        },
        {
            "title":  "Reporting System Issues",
            "fill":   RGBColor(213, 232, 212),
            "border": DARK_GREEN,
            "tcolor": DARK_GREEN,
            "items":  [
                "For technical issues with the PSFGS platform:",
                "- Use the in-app feedback button (bottom right)",
                "- Email: support@glancemanagement.co.za",
                "- Include your municipality name, module affected, and a description of the issue",
                "- Screenshots are helpful — use the in-app screenshot tool",
                "Target response time: 4 business hours",
            ],
        },
        {
            "title":  "Next Steps After Training",
            "fill":   RGBColor(252, 228, 187),
            "border": AMBER,
            "tcolor": RGBColor(140, 80, 5),
            "items":  [
                "1. Complete the post-training assessment (link provided by trainer)",
                "2. Log in to the live system using your credentials",
                "3. Complete the Module 13 Entity Settings setup",
                "4. Add your first real audit finding in AFT",
                "5. Schedule a 30-day review with your System Administrator",
                "6. Share this manual with colleagues who will use the system",
            ],
        },
    ]

    bw = (SLIDE_W - Cm(2.8)) / 3
    bh = SLIDE_H - Cm(4.5)
    by = Cm(3.4)
    for i, box in enumerate(boxes):
        bx = Cm(0.7) + i * (bw + Cm(0.35))
        add_rect(slide, bx, by, bw, bh, fill_rgb=box["fill"], line_rgb=box["border"], line_width_pt=2)
        add_rect(slide, bx, by, bw, Cm(0.9), fill_rgb=box["border"])
        add_textbox(slide, bx + Cm(0.25), by + Cm(0.1), bw - Cm(0.5), Cm(0.75),
                    box["title"], font_size=12, bold=True, color=WHITE)
        txb = slide.shapes.add_textbox(bx + Cm(0.25), by + Cm(1.05), bw - Cm(0.5), bh - Cm(1.2))
        tf  = txb.text_frame
        tf.word_wrap = True
        first = True
        for item in box["items"]:
            if first:
                p = tf.paragraphs[0]; first = False
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(4)
            run = p.add_run()
            run.text = item
            run.font.name  = "Calibri"
            run.font.size  = Pt(10 if item.startswith("-") or item[0].isdigit() else 10.5)
            run.font.bold  = not (item.startswith("-") or item[0].isdigit())
            run.font.color.rgb = BODY

    # Contact footer
    add_rect(slide, Cm(0.7), SLIDE_H - Cm(1.1), SLIDE_W - Cm(1.4), Cm(0.85),
             fill_rgb=DARK_BLUE)
    add_textbox(slide, Cm(1.0), SLIDE_H - Cm(1.0), SLIDE_W - Cm(2.0), Cm(0.75),
                "Glance Management Technologies (Pty) Ltd  |  www.glancemanagement.co.za  |  support@glancemanagement.co.za",
                font_size=10, bold=False, color=WHITE, align=PP_ALIGN.CENTER)


def slide_thankyou(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=DARK_BLUE)
    add_rect(slide, 0, 0, SLIDE_W, Cm(0.6), fill_rgb=MID_BLUE)
    add_rect(slide, 0, SLIDE_H - Cm(0.6), SLIDE_W, Cm(0.6), fill_rgb=MID_BLUE)

    # Decorative center band
    add_rect(slide, 0, SLIDE_H / 2 - Cm(0.08), SLIDE_W, Cm(0.16), fill_rgb=MID_BLUE)

    add_textbox(slide, Cm(1), Cm(2.5), SLIDE_W - Cm(2), Cm(3.5),
                "Thank You", font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, Cm(1), Cm(6.5), SLIDE_W - Cm(2), Cm(1.6),
                "PSFGS — Empowering Clean Governance Across South Africa",
                font_size=18, bold=False, italic=True, color=LIGHT_BLUE2,
                align=PP_ALIGN.CENTER)

    add_rect(slide, Cm(6), Cm(8.5), SLIDE_W - Cm(12), Cm(0.08), fill_rgb=MID_BLUE)

    add_textbox(slide, Cm(1), Cm(9.0), SLIDE_W - Cm(2), Cm(1.2),
                "Glance Management Technologies (Pty) Ltd",
                font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, Cm(1), Cm(10.3), SLIDE_W - Cm(2), Cm(0.9),
                "www.glancemanagement.co.za  |  support@glancemanagement.co.za",
                font_size=11, bold=False, color=LIGHT_BLUE2, align=PP_ALIGN.CENTER)

    add_textbox(slide, Cm(1), Cm(11.5), SLIDE_W - Cm(2), Cm(0.75),
                "Version 1.0  |  April 2026",
                font_size=10, bold=False, color=GREY, align=PP_ALIGN.CENTER)


# ===========================================================================
# MODULE DATA
# ===========================================================================

MODULES = [
    # -------------------------------------------------------------------------
    # 1 — DASH
    # -------------------------------------------------------------------------
    {
        "num":   1,
        "name":  "Executive Dashboard",
        "code":  "DASH",
        "purpose": (
            "The Executive Dashboard provides a consolidated, real-time governance view "
            "of the municipality's financial and compliance health. It serves as the single "
            "pane of glass for senior management, aggregating live data from all 13 modules "
            "into one screen. Colour-coded KPI tiles and AI-generated health analysis allow "
            "leadership to identify issues at a glance without navigating each module individually."
        ),
        "legislation": "MFMA Section 71 (monthly financial reporting) and Section 72 (mid-year assessment)",
        "users": ["Municipal Manager", "Mayor", "CFO", "Councillor"],
        "benefits": [
            "Instant governance health overview — no report-reading required",
            "AI analysis highlights the most urgent issues automatically",
            "Colour-coded RAG status — Red/Amber/Green at a glance",
            "One-click drill-down into any module from the dashboard",
        ],
        "features": [
            ("KPI Tiles", "Live tiles showing status of each module with colour coding"),
            ("AI Health Analysis", "Automated narrative analysis of governance risk areas"),
            ("Module Shortcuts", "One-click navigation to any module directly from tiles"),
            ("Status Indicators", "Colour-coded Red/Amber/Green for immediate risk identification"),
        ],
        "workflow": ["Login to PSFGS", "Dashboard auto-loads", "Review KPI tiles", "Read AI analysis", "Drill down to module"],
        "scenario": (
            "You are the CFO of Buffalo City Municipality. It is the 5th of the month. "
            "Log in to PSFGS and identify which governance area requires urgent attention "
            "based on the dashboard tiles."
        ),
        "tasks": [
            "Login as cfo@buffalocity.gov.za with the test password",
            "Wait for the Executive Dashboard to load automatically",
            "Identify all tiles showing Red (critical) status",
            "Identify all tiles showing Amber (warning) status",
            "Click on the most critical Red tile to navigate to that module",
        ],
        "outcome": (
            "The dashboard displays all 13 module tiles with live status. "
            "You can identify that the Audit Findings Tracker (AFT) has overdue findings "
            "shown in Red, and the IFWE Register has a new item requiring CFO attention in Amber. "
            "The AI analysis panel confirms these as the top two governance risks for the period."
        ),
        "email": "cfo@buffalocity.gov.za",
    },
    # -------------------------------------------------------------------------
    # 2 — AFT
    # -------------------------------------------------------------------------
    {
        "num":   2,
        "name":  "Audit Findings Tracker",
        "code":  "AFT",
        "purpose": (
            "The Audit Findings Tracker provides end-to-end lifecycle management of all "
            "AGSA and internal audit findings, from initial capture through assignment, "
            "implementation, evidence submission and final closure. It eliminates the use "
            "of spreadsheets and ensures every finding is tracked with accountability, "
            "due dates and an audit trail."
        ),
        "legislation": (
            "MFMA Section 131 (remediation of audit findings), Public Audit Act Section 28, "
            "AGSA findings management protocol"
        ),
        "users": ["Auditor", "Municipal Manager", "Director", "Manager", "Officer", "CFO"],
        "benefits": [
            "Full finding lifecycle from capture to closure in one system",
            "AG reference numbers tracked — no finding is ever lost",
            "Automatic overdue escalation alerts to the responsible person",
            "Repeat finding flag prevents the same issue recurring",
        ],
        "features": [
            ("Finding Capture", "Capture all AGSA and IA findings with full detail and AG reference"),
            ("6-Stage Workflow", "Structured workflow: Capture > Assign > In Progress > Evidence > Review > Closed"),
            ("Evidence Submission", "Officers attach evidence documents directly to the finding record"),
            ("Statistics Dashboard", "Live charts showing findings by severity, status and directorate"),
        ],
        "workflow": ["Capture finding", "Assign to director", "Director sets action plan", "Officer submits evidence", "IA reviews", "CFO closes"],
        "scenario": (
            "AGSA has issued a management letter with reference AG-2024-015 regarding "
            "non-compliance with SCM thresholds in the EIS directorate. Capture this finding "
            "and assign it to the EIS Director for remediation."
        ),
        "tasks": [
            "Login as audit@buffalocity.gov.za",
            "Navigate to AFT > New Finding",
            "Complete all fields: AG Ref AG-2024-015, Type: SCM Non-Compliance, Severity: High, Directorate: EIS",
            "Save the finding",
            "Login as cm@buffalocity.gov.za > find the finding > Assign to director.eis@buffalocity.gov.za",
            "Set the due date to 60 days from today and save",
        ],
        "outcome": (
            "The finding AG-2024-015 is captured in the system with status 'Assigned'. "
            "The EIS Director receives a system notification with the finding details and due date. "
            "The finding appears on the dashboard KPI tile for AFT and in the statistics report."
        ),
        "email": "audit@buffalocity.gov.za",
    },
    # -------------------------------------------------------------------------
    # 3 — AVS
    # -------------------------------------------------------------------------
    {
        "num":   3,
        "name":  "Asset Verification System",
        "code":  "AVS",
        "purpose": (
            "The Asset Verification System enables physical verification and condition "
            "assessment of all municipal assets against the asset register. It provides "
            "a structured, digital approach to annual and interim verification sessions, "
            "ensuring that the municipality's asset register remains accurate, GRAP 17 "
            "compliant, and audit-ready at all times."
        ),
        "legislation": (
            "MFMA Section 63 (asset management), GRAP 17 (property, plant and equipment), "
            "Municipal Asset Transfer Regulations"
        ),
        "users": ["Asset Manager", "Director", "Finance Officer"],
        "benefits": [
            "Eliminates paper-based verification forms and manual data capturing",
            "GPS coordinates recorded for each asset in the field",
            "Discrepancy reports generated automatically for CFO review",
            "GRAP 17 condition categories ensure accounting compliance",
        ],
        "features": [
            ("Asset Register", "Full asset register integrated with verification session workflow"),
            ("Condition Scoring", "Standardised scoring: Good / Fair / Poor / Missing / Disposed"),
            ("GPS Capture", "GPS coordinates captured for each physical asset in the field"),
            ("Discrepancy Reports", "Auto-generated reports of missing, poor and unverified assets"),
        ],
        "workflow": ["Create session", "Assign verification team", "Field verification", "Capture condition", "Identify discrepancies", "Report to CFO"],
        "scenario": (
            "Create a new asset verification session for the EIS directorate fleet vehicles. "
            "Verify three assets and mark one vehicle as 'Poor' condition with a note "
            "recommending disposal assessment."
        ),
        "tasks": [
            "Login to AVS module",
            "Select 'New Verification Session' and choose directorate: EIS",
            "Add three vehicle assets from the asset register to the session",
            "Mark the first two as 'Good' condition with verification date",
            "Mark the third vehicle as 'Poor' condition and add a note: 'Recommend disposal assessment'",
            "Save the session and view the auto-generated discrepancy report",
        ],
        "outcome": (
            "The verification session is saved with status 'Completed'. "
            "One vehicle is flagged as 'Poor' condition and appears on the discrepancy report. "
            "The CFO receives a notification that a new discrepancy report is available for review. "
            "The asset register is updated with the latest condition information."
        ),
        "email": "assets@buffalocity.gov.za",
    },
    # -------------------------------------------------------------------------
    # 4 — BVM
    # -------------------------------------------------------------------------
    {
        "num":   4,
        "name":  "Budget vs Actual Monitor",
        "code":  "BVM",
        "purpose": (
            "The Budget vs Actual Monitor provides real-time monitoring of all municipal "
            "budget votes against actual expenditure. It enables CFOs and directors to "
            "identify over-spending and under-spending variances immediately, supporting "
            "proactive financial management and compliance with MFMA Section 71 monthly "
            "reporting requirements."
        ),
        "legislation": (
            "MFMA Section 71 (monthly budget statements), Section 72 (mid-year budget assessment), "
            "Municipal Budget and Reporting Regulations 2009"
        ),
        "users": ["CFO", "Directors", "Finance Officers", "Municipal Manager"],
        "benefits": [
            "Real-time variance visibility without waiting for monthly reports",
            "Automatic alerts when expenditure variance exceeds configurable thresholds",
            "Directorate-level filtering for focused budget management",
            "Trend charts show spending trajectory over the financial year",
        ],
        "features": [
            ("Budget Vote Dashboard", "All votes displayed with budget, actual, variance and percentage"),
            ("Variance Alerts", "Automatic Red/Amber alerts when variances exceed 10% threshold"),
            ("Trend Charts", "Visual spending trend charts per vote and per directorate"),
            ("Period Filter", "Filter by month, quarter or year-to-date for any analysis period"),
        ],
        "workflow": ["View budget votes", "Filter by period and directorate", "Identify variances above 10%", "Flag concern", "CFO review", "Escalate if needed"],
        "scenario": (
            "It is the end of Quarter 2 of the financial year. Review the EIS directorate "
            "budget votes and identify any votes with expenditure variance exceeding 10%. "
            "Prepare to escalate the most critical variances to the Municipal Manager."
        ),
        "tasks": [
            "Login as cfo@buffalocity.gov.za",
            "Navigate to the BVM module",
            "Set the period filter to 'Quarter 2 Year to Date'",
            "Filter the directorate to 'EIS'",
            "Identify all votes with variance percentage exceeding 10% (shown in red)",
            "Note the vote numbers and variance amounts for escalation",
        ],
        "outcome": (
            "The BVM dashboard displays all EIS votes for Q2. At least two votes show "
            "variance exceeding 10% and are highlighted in red. The trend chart shows "
            "the spending trajectory per month. You can export this view as a PDF "
            "for the Monthly Budget Statement report."
        ),
        "email": "cfo@buffalocity.gov.za",
    },
    # -------------------------------------------------------------------------
    # 5 — IFW
    # -------------------------------------------------------------------------
    {
        "num":   5,
        "name":  "IFWE Register",
        "code":  "IFW",
        "purpose": (
            "The IFWE Register captures and tracks all Irregular, Fruitless and Wasteful "
            "Expenditure identified during the financial year. It provides a complete audit "
            "trail from identification through investigation to condonation or recovery, "
            "and automatically populates the AFS disclosure note and Council reports."
        ),
        "legislation": (
            "MFMA Section 32 (irregular expenditure), Section 34 (fruitless and wasteful expenditure), "
            "Section 38 (financial misconduct), Treasury Regulation 9.1"
        ),
        "users": ["CFO", "Finance Director", "Internal Audit", "Municipal Manager"],
        "benefits": [
            "Single register for all three IFWE categories — no separate spreadsheets",
            "AFS disclosure note populated automatically from register data",
            "Investigation status tracked with responsible officer and due dates",
            "Council reporting made easy with one-click consolidated report",
        ],
        "features": [
            ("IFWE Classification", "Classify each item as Irregular, Fruitless or Wasteful with full detail"),
            ("Investigation Tracking", "Track investigation status, responsible officer and outcome"),
            ("Condonation Register", "Record Council condonation decisions with resolution references"),
            ("AFS Disclosure Feed", "Automatically feeds the AFS IFWE disclosure note for year-end"),
        ],
        "workflow": ["Identify IFWE", "Capture in register", "Investigate", "Determine responsibility", "Condonation or recovery", "Council disclosure", "AFS note"],
        "scenario": (
            "A payment of R450,000 was made to a supplier without following the three-quotation "
            "process required by SCM policy. Capture this as Irregular Expenditure and set its "
            "investigation status."
        ),
        "tasks": [
            "Login to the IFW module",
            "Select 'New IFWE Item'",
            "Classify as 'Irregular Expenditure'",
            "Enter amount: R450,000 and payment date",
            "Describe the non-compliance: 'Payment made without three quotations as required by SCM policy'",
            "Set investigation status to 'Under Investigation' and assign to Internal Audit",
            "Save the record",
        ],
        "outcome": (
            "The item is captured in the IFWE Register as Irregular Expenditure with status "
            "'Under Investigation'. It appears on the dashboard IFWE tile and in the "
            "year-to-date irregular expenditure total. Internal Audit is notified. "
            "The item will auto-populate in the AFS IFWE disclosure note at year-end."
        ),
        "email": "cfo@buffalocity.gov.za",
    },
    # -------------------------------------------------------------------------
    # 6 — SCC
    # -------------------------------------------------------------------------
    {
        "num":   6,
        "name":  "SCM Compliance Checker",
        "code":  "SCC",
        "purpose": (
            "The SCM Compliance Checker automates the compliance checking of procurement "
            "requisitions against SCM policy and legislation. It runs a 12-point rule engine "
            "on each requisition and produces a compliance score from 0 to 100 percent, "
            "highlighting specific exceptions that must be resolved before the procurement "
            "can proceed."
        ),
        "legislation": (
            "MFMA SCM Regulations 2005, Preferential Procurement Policy Framework Act (PPPFA), "
            "National Treasury Instruction Notes 3 and 7 of 2021/22"
        ),
        "users": ["SCM Director", "Demand Manager", "CFO", "Internal Audit"],
        "benefits": [
            "Prevents irregular expenditure before it happens — not after",
            "12-point automated rule engine covers all common compliance failures",
            "Compliance score gives instant visibility of procurement health",
            "Exceptions report supports accountability and corrective action",
        ],
        "features": [
            ("Requisition Capture", "Capture all procurement requisitions with supplier and amount detail"),
            ("12-Point Rule Engine", "Automated compliance check against SCM policy, thresholds and NT rules"),
            ("Tax Clearance Check", "Validates supplier tax clearance status before approval"),
            ("Compliance Score", "0-100% compliance score with colour coding and exception detail"),
        ],
        "workflow": ["Capture requisition", "Run compliance check", "Review score and exceptions", "Fix exceptions", "Resubmit", "Approve"],
        "scenario": (
            "Capture a new procurement requisition for the supply of office furniture at "
            "R185,000. Run the automated compliance check and review the compliance score "
            "and any exceptions that are identified."
        ),
        "tasks": [
            "Login to the SCC module",
            "Select 'New Requisition'",
            "Complete all fields: description, amount R185,000, supplier name, category",
            "Save the requisition",
            "Click 'Run Compliance Check'",
            "Review the compliance score and read each exception listed",
            "Note which exceptions need to be resolved before approval",
        ],
        "outcome": (
            "The compliance check completes and shows a score (e.g., 75%). "
            "Exceptions are listed with descriptions — for example, missing tax clearance "
            "certificate or incorrect quotation category. You can resolve exceptions, "
            "rerun the check, and resubmit once the score reaches 100%."
        ),
        "email": "scm@buffalocity.gov.za",
    },
    # -------------------------------------------------------------------------
    # 7 — CMT
    # -------------------------------------------------------------------------
    {
        "num":   7,
        "name":  "Consequence Management Tracker",
        "code":  "CMT",
        "purpose": (
            "The Consequence Management Tracker provides structured tracking of all financial "
            "misconduct cases, disciplinary proceedings and AGSA referrals. It ensures that "
            "the municipality complies with the COGTA Consequence Management Framework "
            "and that every case is progressed, documented and reported to Council."
        ),
        "legislation": (
            "MFMA Section 171-173 (financial misconduct proceedings), "
            "Municipal Systems Act Section 67 (disciplinary action), "
            "COGTA Consequence Management Framework"
        ),
        "users": ["Municipal Manager", "HR Director", "Internal Audit", "Legal Advisor"],
        "benefits": [
            "AGSA referrals tracked alongside internal disciplinary cases in one register",
            "Hearing dates, outcomes and appeal decisions all captured digitally",
            "Council reporting on consequence management is automated",
            "Demonstrates accountability to AGSA and oversight bodies",
        ],
        "features": [
            ("Case Capture", "Capture all financial misconduct and disciplinary cases with full detail"),
            ("Case Type Classification", "Classify as Financial Misconduct, Criminal Referral or Civil Recovery"),
            ("Hearing Management", "Schedule hearings, record outcomes and manage appeals"),
            ("AGSA Referral Tracking", "Track cases referred to or by the Auditor-General"),
        ],
        "workflow": ["Open case", "Assign investigator", "Schedule hearing", "Capture outcome", "Record appeal", "Close case", "Report to Council"],
        "scenario": (
            "Open a consequence management case against a SCM official who awarded "
            "a contract of R2.3 million without following the required competitive bidding "
            "process, resulting in Irregular Expenditure being raised."
        ),
        "tasks": [
            "Login to the CMT module",
            "Select 'New Case'",
            "Capture the respondent's name, employee number and directorate",
            "Set the case type to 'Financial Misconduct'",
            "Describe the allegation: 'Contract awarded without competitive bidding — R2.3m'",
            "Assign an investigating officer",
            "Set the expected hearing date and save the case",
        ],
        "outcome": (
            "The case is created with status 'Under Investigation'. "
            "The investigating officer is notified by the system. "
            "The case appears on the Consequence Management dashboard tile "
            "and will be included in the next Council report on consequence management."
        ),
        "email": "cm@buffalocity.gov.za",
    },
    # -------------------------------------------------------------------------
    # 8 — DAR
    # -------------------------------------------------------------------------
    {
        "num":   8,
        "name":  "Delegation of Authority Register",
        "code":  "DAR",
        "purpose": (
            "The Delegation of Authority Register provides a digital register of all Council "
            "delegations and sub-delegations with associated financial approval limits. It "
            "replaces paper-based delegation files with a searchable, auditable system that "
            "tracks enabling resolutions, validity periods and expiry alerts."
        ),
        "legislation": (
            "MFMA Section 79 (delegation of powers), Municipal Systems Act Section 59, "
            "Municipal Finance Management Act general delegation provisions"
        ),
        "users": ["Municipal Manager", "CFO", "Directors", "Company Secretary"],
        "benefits": [
            "Digital register eliminates the risk of using outdated delegation documents",
            "Financial approval limits clearly linked to each delegatee",
            "Expiry alerts prevent actions under lapsed delegations",
            "Council resolution references provide legal traceability",
        ],
        "features": [
            ("Delegation Hierarchy", "Full hierarchy from Council to MM to CFO to Directors"),
            ("Financial Limits", "Monetary approval thresholds per delegatee clearly recorded"),
            ("Resolution Linking", "Each delegation linked to enabling Council resolution number"),
            ("Expiry Alerts", "Automatic alerts 30 days before any delegation expires"),
        ],
        "workflow": ["Council passes resolution", "Capture delegation in system", "Assign delegatee", "Set financial limit", "Set expiry date", "Annual review"],
        "scenario": (
            "Capture a delegation from the Municipal Manager to the SCM Director "
            "for contract approvals up to R500,000, valid for the current financial year. "
            "Link it to Council Resolution 45/2025."
        ),
        "tasks": [
            "Login to the DAR module",
            "Select 'New Delegation'",
            "Select the delegator: Municipal Manager",
            "Select the delegatee: SCM Director",
            "Set the financial limit: R500,000",
            "Link the enabling resolution: Council Resolution 45/2025",
            "Set the validity period: 1 July 2025 to 30 June 2026",
            "Save the delegation",
        ],
        "outcome": (
            "The delegation is saved in the register with status 'Active'. "
            "The SCM Director can view the delegation from their profile. "
            "The system will send an alert 30 days before 30 June 2026 when the delegation expires. "
            "The delegation is available for audit inspection in the DAR module."
        ),
        "email": "cm@buffalocity.gov.za",
    },
    # -------------------------------------------------------------------------
    # 9 — PCR
    # -------------------------------------------------------------------------
    {
        "num":   9,
        "name":  "Policy Compliance Manager",
        "code":  "PCR",
        "purpose": (
            "The Policy Compliance Manager tracks all municipal policies, their review "
            "schedules and compliance against the full list of MFMA-required policies. "
            "It provides a gap analysis between policies that are required and those "
            "that are in place and up to date, with automatic 30-day review alerts."
        ),
        "legislation": (
            "MFMA various sections (Budget Policy, SCM Policy, Treasury Policy, etc.), "
            "COGTA Policy Framework, National Treasury policy guidelines"
        ),
        "users": ["CFO", "Directors", "Company Secretary", "Internal Audit"],
        "benefits": [
            "Full policy register with required vs in-place gap analysis",
            "Automatic 30-day alert before any policy review date",
            "Council approval tracking ensures policies are formally adopted",
            "AGSA no longer finds outdated or missing policies",
        ],
        "features": [
            ("Policy Register", "Complete register of all municipal policies by category"),
            ("Review Schedule", "Automatic review date tracking with 30-day advance alerts"),
            ("Gap Analysis", "Required vs in-place comparison showing policy coverage"),
            ("Owner Assignment", "Each policy assigned to a responsible directorate owner"),
        ],
        "workflow": ["Capture policy", "Assign owner", "Set review date", "Receive 30-day alert", "Review and update", "Council reapproval", "Update status"],
        "scenario": (
            "The SCM Policy is due for its annual review this month. Update the policy "
            "status to 'Under Review', assign the review to the SCM Director, "
            "and set a new review completion target date."
        ),
        "tasks": [
            "Login to the PCR module",
            "Find the SCM Policy in the policy register",
            "Click 'Edit Policy'",
            "Change the status from 'Current' to 'Under Review'",
            "Assign the review to the SCM Director",
            "Set the review completion target date to 30 days from today",
            "Save the changes",
        ],
        "outcome": (
            "The SCM Policy status is updated to 'Under Review'. "
            "The SCM Director is notified by the system. "
            "The policy appears on the review tracking dashboard. "
            "Once the review is completed, the status will be updated to 'Current' "
            "after Council approval is captured."
        ),
        "email": "cfo@buffalocity.gov.za",
    },
    # -------------------------------------------------------------------------
    # 10 — S71
    # -------------------------------------------------------------------------
    {
        "num":   10,
        "name":  "Section 71 Auto-Generator",
        "code":  "S71",
        "purpose": (
            "The Section 71 Auto-Generator automatically compiles the monthly Section 71 "
            "financial report for submission to National Treasury, drawing data directly "
            "from the BVM module. It produces reports in the NT-prescribed format with "
            "prior-period comparisons and supports the CFO sign-off workflow."
        ),
        "legislation": (
            "MFMA Section 71 (monthly budget statements to National Treasury), "
            "Municipal Budget and Reporting Regulations 2009 — prescribed format and 10-day deadline"
        ),
        "users": ["CFO", "Finance Director", "Municipal Manager"],
        "benefits": [
            "Eliminates manual report compilation — saves 8+ hours per month",
            "NT-prescribed format guaranteed — no formatting errors",
            "Prior period comparison auto-calculated from BVM data",
            "Late submission alerts ensure the 10th-of-month deadline is never missed",
        ],
        "features": [
            ("Auto-Compilation", "Pulls BVM data automatically — no manual data entry required"),
            ("NT-Prescribed Format", "Output exactly matches National Treasury reporting template"),
            ("CFO Sign-Off Workflow", "Digital CFO approval before report is finalised"),
            ("Submission Tracking", "Records submission date and tracks late submissions"),
        ],
        "workflow": ["Month-end trigger", "System pulls BVM data", "CFO reviews report", "CFO makes corrections", "CFO approves", "PDF generated", "Submit to NT by 10th"],
        "scenario": (
            "It is the 8th of the month. Generate the Section 71 report for the previous "
            "month, review the auto-compiled data, and approve it for submission to "
            "National Treasury before the 10th deadline."
        ),
        "tasks": [
            "Login as cfo@buffalocity.gov.za",
            "Navigate to the S71 module",
            "Select the reporting period (previous month)",
            "Review the auto-compiled Section 71 report",
            "Check the figures against the BVM for accuracy",
            "Click 'Approve' to apply the CFO digital sign-off",
            "Export the final PDF for submission to National Treasury",
        ],
        "outcome": (
            "The Section 71 report is approved and available as a PDF. "
            "The system records the submission date and CFO approval. "
            "The report appears in the submission tracking register. "
            "If submitted after the 10th, the system flags it as a late submission "
            "for audit trail purposes."
        ),
        "email": "cfo@buffalocity.gov.za",
    },
    # -------------------------------------------------------------------------
    # 11 — PKD
    # -------------------------------------------------------------------------
    {
        "num":   11,
        "name":  "Performance KPI Dashboard",
        "code":  "PKD",
        "purpose": (
            "The Performance KPI Dashboard monitors all IDP-linked SDBIP key performance "
            "indicators by quarter. It enables directors to submit actual performance against "
            "quarterly targets, and generates RAG status (Red/Amber/Green) for each KPI, "
            "supporting the Section 46 annual performance report."
        ),
        "legislation": (
            "Municipal Systems Act Section 41 (performance management), "
            "Municipal Planning and Performance Management Regulations 2001, "
            "Section 46 annual performance report requirements"
        ),
        "users": ["Municipal Manager", "Directors", "Performance Manager", "Council"],
        "benefits": [
            "SDBIP KPIs managed in one place — no multiple spreadsheets",
            "RAG status calculated automatically from targets vs actuals",
            "Quarterly scorecards ready for Section 71 and Section 46 reporting",
            "Corrective action plans can be captured directly against underperforming KPIs",
        ],
        "features": [
            ("KPI Capture", "Capture all SDBIP KPIs with quarterly targets and unit of measure"),
            ("Actual Submission", "Directors submit quarterly actual performance against each KPI"),
            ("RAG Calculation", "Automatic Red/Amber/Green status based on target achievement"),
            ("Directorate Scorecard", "Per-directorate scorecard report for management review"),
        ],
        "workflow": ["Capture KPIs and targets", "Quarter end", "Directors submit actuals", "System calculates RAG", "Review underperformers", "Capture corrective actions"],
        "scenario": (
            "Quarter 2 of the financial year has ended. Capture the actual performance "
            "for three EIS directorate KPIs and identify which ones are underperforming "
            "against their targets."
        ),
        "tasks": [
            "Login to the PKD module",
            "Select directorate: EIS and period: Quarter 2",
            "Open the first KPI and enter the Q2 actual value",
            "Open the second KPI and enter the Q2 actual value",
            "Open the third KPI and enter the Q2 actual value",
            "Review the RAG status calculated for each KPI",
            "Note the Red KPIs and what corrective actions are needed",
        ],
        "outcome": (
            "The three KPI actuals are saved. The system calculates the RAG status: "
            "one KPI shows Green (achieved), one shows Amber (within 10% of target), "
            "and one shows Red (below target by more than 10%). "
            "The directorate scorecard updates automatically and is available for the "
            "next EXCO meeting report."
        ),
        "email": "performance@buffalocity.gov.za",
    },
    # -------------------------------------------------------------------------
    # 12 — ERP
    # -------------------------------------------------------------------------
    {
        "num":   12,
        "name":  "Executive Reporting Portal",
        "code":  "ERP",
        "purpose": (
            "The Executive Reporting Portal consolidates governance and financial data from "
            "all 13 modules into committee-ready report packs for the Mayor, Municipal Manager "
            "and Council committees. It includes AI-generated narrative analysis and allows "
            "period selection, PDF export and trend comparison."
        ),
        "legislation": (
            "MFMA Section 129 (annual report), MSA Section 46 (annual performance report), "
            "MFMA Section 71/72 (monthly and mid-year reports)"
        ),
        "users": ["Mayor", "Municipal Manager", "CFO", "Council", "MPAC"],
        "benefits": [
            "All 13 modules summarised in one executive report — no manual compilation",
            "AI-generated narrative saves hours of report writing time",
            "PDF export ready for distribution to Council and oversight bodies",
            "Trend comparison across periods shows governance improvement over time",
        ],
        "features": [
            ("Consolidated Dashboard", "Pulls live data from all 13 modules into one governance view"),
            ("AI Narrative Analysis", "Automated narrative highlights key risks and achievements"),
            ("Committee Report Packs", "One-click PDF export of committee-ready report packages"),
            ("Period Comparison", "Compare current period with prior quarter or prior year"),
        ],
        "workflow": ["Select reporting period", "System compiles from all modules", "Review AI narrative", "Edit as needed", "Export PDF", "Submit to committee"],
        "scenario": (
            "Prepare the governance report for the next MPAC committee meeting "
            "covering Quarter 2 of the current financial year. The report must include "
            "financial performance, audit findings status and KPI achievement."
        ),
        "tasks": [
            "Login to the ERP module",
            "Select the reporting period: Quarter 2",
            "Allow the system to compile the consolidated governance report",
            "Review the AI-generated narrative analysis",
            "Check that all module data sections are complete",
            "Export the report as a PDF for distribution to MPAC members",
        ],
        "outcome": (
            "A complete Q2 governance report is generated as a PDF. "
            "It includes budget performance summary (from BVM), audit findings status (from AFT), "
            "KPI achievement (from PKD), IFWE totals (from IFW) and consequence management update (from CMT). "
            "The AI narrative highlights the top three governance risks for the quarter."
        ),
        "email": "mayor@buffalocity.gov.za",
    },
    # -------------------------------------------------------------------------
    # 13 — SET
    # -------------------------------------------------------------------------
    {
        "num":   13,
        "name":  "Entity Settings",
        "code":  "SET",
        "purpose": (
            "Entity Settings is the administrative configuration module that allows the "
            "System Administrator to set up the municipality's profile in PSFGS. "
            "This includes the entity name, logo, demarcation code, legislation type, "
            "leadership names and contact details. These settings appear throughout the "
            "system and on all exported reports."
        ),
        "legislation": "Administrative — no specific legislation. Supports all MFMA reporting requirements.",
        "users": ["System_Admin"],
        "benefits": [
            "Municipality branding appears on all system outputs and reports",
            "Correct entity details prevent errors on Section 71 and AFS disclosures",
            "Leadership names auto-populated on signature blocks in reports",
            "Financial year end month setting drives all reporting period logic",
        ],
        "features": [
            ("Entity Profile", "Municipality name, abbreviation, demarcation code and province"),
            ("Logo Upload", "Upload municipal crest (PNG/JPG) — appears on login page and reports"),
            ("Leadership Setup", "Capture names of MM, CFO, Mayor, Speaker and Chief Audit Executive"),
            ("Financial Year Config", "Set financial year end month to drive all reporting periods"),
        ],
        "workflow": ["Login as System Admin", "Open Settings module", "Upload logo", "Complete all tabs", "Save settings", "Verify logo on login page"],
        "scenario": (
            "Configure the PSFGS system for your municipality. Upload the municipal "
            "crest as the logo, complete the General information tab, the Contact details "
            "tab and the Leadership names tab."
        ),
        "tasks": [
            "Login as admin@buffalocity.gov.za",
            "Navigate to Settings > Entity Settings",
            "Upload the municipal crest PNG file as the logo",
            "Complete the General tab: entity name, abbreviation, demarcation code, province",
            "Complete the Contact tab: physical address, postal address, telephone and website",
            "Complete the Leadership tab: MM name, CFO name, Mayor name, Speaker name",
            "Set the financial year end month to June",
            "Click 'Save All Settings'",
        ],
        "outcome": (
            "The entity settings are saved. The municipal logo appears on the PSFGS login page "
            "and in the sidebar navigation. All exported reports display the municipality's "
            "name and crest. Leadership names auto-populate on report signature blocks. "
            "The financial year is set to July to June."
        ),
        "email": "admin@buffalocity.gov.za",
    },
]


# ===========================================================================
# MAIN
# ===========================================================================

def main():
    prs = new_prs()

    # ---- Foundation slides ----
    slide_title(prs)
    slide_agenda(prs)
    slide_overview(prs)
    slide_login(prs)
    slide_navigation(prs)
    slide_user_roles(prs)

    # ---- 13 × 3 module slides ----
    for mod in MODULES:
        module_intro(
            prs,
            mod_num       = mod["num"],
            mod_name      = mod["name"],
            mod_code      = mod["code"],
            purpose_text  = mod["purpose"],
            legislation   = mod["legislation"],
            users_list    = mod["users"],
            key_benefits  = mod["benefits"],
        )
        module_features(
            prs,
            mod_name       = mod["name"],
            features_4     = mod["features"],
            workflow_steps = mod["workflow"],
        )
        module_exercise(
            prs,
            mod_name        = mod["name"],
            scenario_text   = mod["scenario"],
            task_steps      = mod["tasks"],
            expected_outcome= mod["outcome"],
            test_email      = mod["email"],
        )

    # ---- Appendix slides ----
    slide_roles_matrix(prs)
    slide_glossary(prs)
    slide_support(prs)
    slide_thankyou(prs)

    prs.save(OUTPUT_PATH)
    print(f"SUCCESS: Presentation saved to:\n{OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
